"""
pgxllm PPT 생성 스크립트
실행: python docs/make_ppt.py
출력: docs/pgxllm.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ───────────────────────────────────────────────
NAVY   = RGBColor(0x1e, 0x3a, 0x5f)   # 사이드바 색
TEAL   = RGBColor(0x0d, 0x9e, 0x8a)   # 포인트 색
WHITE  = RGBColor(0xff, 0xff, 0xff)
GRAY   = RGBColor(0x6b, 0x72, 0x80)
LGRAY  = RGBColor(0xf3, 0xf4, 0xf6)
DGRAY  = RGBColor(0x37, 0x41, 0x51)
ORANGE = RGBColor(0xf5, 0x9e, 0x0b)
GREEN  = RGBColor(0x10, 0xb9, 0x81)
RED    = RGBColor(0xef, 0x44, 0x44)
MINT   = RGBColor(0xe6, 0xfb, 0xf7)

W = Inches(13.33)   # 와이드 슬라이드 너비
H = Inches(7.5)     # 높이

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ── 헬퍼 함수 ─────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=Pt(14), bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def slide_header(slide, title, subtitle=None, dark=True):
    """상단 헤더 바."""
    bg = NAVY if dark else TEAL
    add_rect(slide, 0, 0, W, Inches(1.1), fill=bg)
    add_text(slide, title,
             Inches(0.4), Inches(0.18), Inches(10), Inches(0.55),
             size=Pt(24), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.72), Inches(10), Inches(0.35),
                 size=Pt(13), color=RGBColor(0xb2, 0xf5, 0xea))


def footer(slide, page_n, total):
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill=LGRAY)
    add_text(slide, "pgxllm — PostgreSQL Text-to-SQL System",
             Inches(0.3), H - Inches(0.33), Inches(9), Inches(0.3),
             size=Pt(9), color=GRAY)
    add_text(slide, f"{page_n} / {total}",
             W - Inches(1.2), H - Inches(0.33), Inches(1.0), Inches(0.3),
             size=Pt(9), color=GRAY, align=PP_ALIGN.RIGHT)


def bullet_block(slide, items, x, y, w,
                 bullet="•", size=Pt(13), color=DGRAY, gap=Inches(0.38)):
    """간단 불릿 목록."""
    cy = y
    for item in items:
        add_text(slide, f"{bullet}  {item}", x, cy, w, gap,
                 size=size, color=color)
        cy += gap
    return cy


def card(slide, x, y, w, h, title, body_lines,
         hdr_fill=TEAL, hdr_color=WHITE,
         body_fill=WHITE, body_color=DGRAY,
         title_size=Pt(13), body_size=Pt(11.5)):
    """제목+본문 카드."""
    hh = Inches(0.42)
    add_rect(slide, x, y, w, hh, fill=hdr_fill)
    add_text(slide, title, x + Inches(0.12), y + Inches(0.06),
             w - Inches(0.15), hh - Inches(0.06),
             size=title_size, bold=True, color=hdr_color)
    bh = h - hh
    add_rect(slide, x, y + hh, w, bh, fill=body_fill,
             line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(0.75))
    cy = y + hh + Inches(0.12)
    for line in body_lines:
        add_text(slide, line, x + Inches(0.15), cy,
                 w - Inches(0.2), Inches(0.32),
                 size=body_size, color=body_color)
        cy += Inches(0.32)


def table_rows(slide, headers, rows, x, y, w,
               col_widths=None, hdr_fill=NAVY, row_alt=LGRAY):
    """간단 테이블."""
    n_col = len(headers)
    if col_widths is None:
        col_widths = [w / n_col] * n_col
    rh = Inches(0.36)

    # 헤더
    cx = x
    for i, h_txt in enumerate(headers):
        add_rect(slide, cx, y, col_widths[i], rh, fill=hdr_fill)
        add_text(slide, h_txt, cx + Inches(0.08), y + Inches(0.06),
                 col_widths[i] - Inches(0.1), rh - Inches(0.06),
                 size=Pt(11), bold=True, color=WHITE)
        cx += col_widths[i]

    # 행
    for ri, row in enumerate(rows):
        cy = y + rh * (ri + 1)
        fill = LGRAY if ri % 2 == 0 else WHITE
        cx   = x
        for ci, cell in enumerate(row):
            add_rect(slide, cx, cy, col_widths[ci], rh, fill=fill,
                     line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(0.5))
            add_text(slide, cell, cx + Inches(0.08), cy + Inches(0.06),
                     col_widths[ci] - Inches(0.1), rh - Inches(0.06),
                     size=Pt(10.5), color=DGRAY)
            cx += col_widths[ci]


def flow_box(slide, label, x, y, w=Inches(2.0), h=Inches(0.65),
             fill=TEAL, color=WHITE, size=Pt(12)):
    add_rect(slide, x, y, w, h, fill=fill)
    add_text(slide, label, x, y, w, h, size=size, bold=True,
             color=color, align=PP_ALIGN.CENTER)


def arrow_down(slide, x, y, length=Inches(0.35)):
    from pptx.util import Pt
    ln = slide.shapes.add_connector(1,
        x, y, x, y + length)
    ln.line.color.rgb = GRAY
    ln.line.width = Pt(1.5)


# ═══════════════════════════════════════════════════════════════
# 슬라이드 정의
# ═══════════════════════════════════════════════════════════════
TOTAL = 12   # 총 슬라이드 수

# ────────────────────────────────────────────────────────────────
# 1. 표지
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)

# 배경 분할
add_rect(sl, 0, 0, W * 0.42, H, fill=NAVY)
add_rect(sl, W * 0.42, 0, W * 0.58, H, fill=WHITE)

# 좌측 텍스트
add_text(sl, "pgxllm",
         Inches(0.5), Inches(1.8), Inches(5), Inches(1.1),
         size=Pt(52), bold=True, color=WHITE)
add_text(sl, "PostgreSQL Text-to-SQL System",
         Inches(0.5), Inches(2.9), Inches(5), Inches(0.6),
         size=Pt(18), color=TEAL)
add_text(sl, "자연어로 묻고, SQL로 답하다",
         Inches(0.5), Inches(3.55), Inches(5), Inches(0.5),
         size=Pt(14), color=RGBColor(0xb2, 0xf5, 0xea), italic=True)

# 구분선
add_rect(sl, Inches(0.5), Inches(4.15), Inches(4.5), Inches(0.04), fill=TEAL)

# 핵심 키워드
keywords = ["Multi-LLM Support", "Schema Linking", "Query Tuning",
            "Explain Visualize", "Dialect Rules", "Semantic Cache"]
cy = Inches(4.4)
for kw in keywords:
    add_text(sl, f"▸  {kw}", Inches(0.6), cy, Inches(4.5), Inches(0.38),
             size=Pt(12), color=WHITE)
    cy += Inches(0.38)

# 우측 아키텍처 요약
add_text(sl, "System Flow",
         Inches(6.0), Inches(1.2), Inches(6.5), Inches(0.5),
         size=Pt(16), bold=True, color=NAVY)

steps = [
    ("자연어 질문",      TEAL),
    ("S1  질문 이해",    NAVY),
    ("S2  스키마 링킹",  NAVY),
    ("S3  SQL 생성 (LLM)", NAVY),
    ("S4  검증 · 보정",  NAVY),
    ("SQL + 설명 반환",  RGBColor(0x10, 0xb9, 0x81)),
]
sy = Inches(1.8)
for label, col in steps:
    add_rect(sl, Inches(7.2), sy, Inches(4.5), Inches(0.6), fill=col)
    add_text(sl, label, Inches(7.2), sy, Inches(4.5), Inches(0.6),
             size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if label != steps[-1][0]:
        add_rect(sl, Inches(9.2), sy + Inches(0.6), Inches(0.5), Inches(0.22),
                 fill=GRAY)
    sy += Inches(0.82)

footer(sl, 1, TOTAL)


# ────────────────────────────────────────────────────────────────
# 2. 핵심 기능 개요
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "핵심 기능 개요", "pgxllm이 제공하는 8가지 핵심 기능")

features = [
    ("Text-to-SQL",      TEAL,   "자연어 질문 → SQL 자동 변환\n4단계 파이프라인 (S1~S4)"),
    ("다중 LLM 지원",     NAVY,   "Ollama · vLLM · OpenAI\nAnthropic · IBM watsonx.ai"),
    ("스키마 탐색",       TEAL,   "pg_catalog 기반 테이블·컬럼\n인덱스·통계 실시간 조회"),
    ("쿼리 분석",         NAVY,   "pg_stat_statements 수집\nEXPLAIN ANALYZE 시각화"),
    ("쿼리 튜닝",         TEAL,   "LLM 기반 최적화 제안\n인덱스 추가 · 재작성 권고"),
    ("Graph 관계",        NAVY,   "FK · 분석 · 추론 기반\n테이블 관계 그래프 관리"),
    ("Dialect Rules",    TEAL,   "컬럼별 SQL 작성 규칙\n자동 감지 + 수동 등록"),
    ("Semantic Cache",   NAVY,   "TF-IDF 유사 질문 캐시\nLLM 호출 최소화"),
]

cols = 4
rows_data = [features[:4], features[4:]]
for ri, row_items in enumerate(rows_data):
    for ci, (title, col, body) in enumerate(row_items):
        x = Inches(0.3) + ci * Inches(3.22)
        y = Inches(1.35) + ri * Inches(2.7)
        w = Inches(3.05)
        h = Inches(2.45)
        add_rect(sl, x, y, w, h, fill=col)
        # 아이콘 영역
        add_rect(sl, x, y, w, Inches(0.55), fill=RGBColor(0,0,0))  # 투명 효과 대신
        sh = sl.shapes.add_shape(1, x, y, w, Inches(0.55))
        sh.fill.solid(); sh.fill.fore_color.rgb = col
        sh.line.fill.background()
        add_text(sl, title, x + Inches(0.15), y + Inches(0.1),
                 w - Inches(0.2), Inches(0.42),
                 size=Pt(14), bold=True, color=WHITE)
        add_text(sl, body, x + Inches(0.15), y + Inches(0.65),
                 w - Inches(0.25), Inches(1.7),
                 size=Pt(12), color=WHITE)

footer(sl, 2, TOTAL)


# ────────────────────────────────────────────────────────────────
# 3. 시스템 아키텍처
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "시스템 아키텍처", "컴포넌트 구성 및 데이터 흐름")

# 좌측: 컴포넌트 트리
add_text(sl, "컴포넌트 구성", Inches(0.3), Inches(1.25), Inches(5.5), Inches(0.4),
         size=Pt(13), bold=True, color=NAVY)

components = [
    ("core/",            TEAL,   True),
    ("  llm/factory.py", DGRAY,  False),
    ("  pipeline.py",    DGRAY,  False),
    ("  s1~s4 단계",     DGRAY,  False),
    ("intelligence/",    NAVY,   True),
    ("  refresh.py",     DGRAY,  False),
    ("  rule_engine.py", DGRAY,  False),
    ("graph/",           TEAL,   True),
    ("  postgresql / age / neo4j", DGRAY, False),
    ("cache/  (TF-IDF)", NAVY,   True),
    ("parser/ (ANTLR4)", NAVY,   True),
    ("web/app.py",       TEAL,   True),
    ("cli.py",           DGRAY,  True),
]
cy = Inches(1.72)
for comp, col, bold in components:
    add_text(sl, comp, Inches(0.4), cy, Inches(5.2), Inches(0.32),
             size=Pt(11), bold=bold, color=col)
    cy += Inches(0.32)

# 우측: 파이프라인 흐름도
add_text(sl, "Core Pipeline 흐름", Inches(6.2), Inches(1.25), Inches(6.5), Inches(0.4),
         size=Pt(13), bold=True, color=NAVY)

pipeline_steps = [
    ("자연어 질문 입력",          TEAL),
    ("Semantic Cache 확인",      RGBColor(0x64, 0x74, 0x8b)),
    ("S1  질문 이해 · 패턴 감지", NAVY),
    ("S2  스키마 링킹 · JOIN 경로", NAVY),
    ("S3  LLM SQL 생성",          TEAL),
    ("S4  SQL 검증 · 자동 보정",   NAVY),
    ("결과 캐시 저장 → SQL 반환", GREEN),
]
sx = Inches(7.8)
sy = Inches(1.72)
bw = Inches(4.2)
bh = Inches(0.55)
gap = Inches(0.18)

for i, (label, col) in enumerate(pipeline_steps):
    add_rect(sl, sx, sy, bw, bh, fill=col)
    add_text(sl, label, sx, sy, bw, bh,
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipeline_steps) - 1:
        add_rect(sl, sx + bw/2 - Inches(0.06), sy + bh,
                 Inches(0.12), gap, fill=GRAY)
    sy += bh + gap

# Cache HIT 분기 표시
add_text(sl, "↩ Cache HIT → 즉시 반환",
         Inches(6.0), Inches(2.26), Inches(1.7), Inches(0.4),
         size=Pt(9.5), color=ORANGE, italic=True)

footer(sl, 3, TOTAL)


# ────────────────────────────────────────────────────────────────
# 4. Core Pipeline 상세
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "Core Pipeline 상세", "4단계 Text-to-SQL 파이프라인")

stages = [
    ("S1", "질문 이해",
     ["DynamicPattern 감지", "(TOP-N / GROUP BY 등)", "키워드 추출", "후보 테이블 검색", "(pg_trgm fulltext)"],
     TEAL),
    ("S2", "스키마 링킹",
     ["후보 테이블 컬럼 상세", "로드", "Graph JOIN 경로 탐색", "Dialect Rules 수집", "LinkedSchema 구성"],
     NAVY),
    ("S3", "SQL 생성 (LLM)",
     ["System Prompt 조립", "(규칙+패턴+Few-shot)", "User Prompt 조립", "(스키마+질문)", "LLM 호출 → SQL 파싱"],
     TEAL),
    ("S4", "SQL 검증·보정",
     ["EXPLAIN 문법 검증", "오류 시 LLM 재생성", "(최대 3회)", "성공 시 캐시 저장", "최종 SQL 반환"],
     NAVY),
]

for i, (num, title, bullets, col) in enumerate(stages):
    x = Inches(0.25) + i * Inches(3.25)
    y = Inches(1.25)
    w = Inches(3.1)

    # 헤더
    add_rect(sl, x, y, w, Inches(0.9), fill=col)
    add_text(sl, num, x + Inches(0.1), y + Inches(0.05), Inches(0.55), Inches(0.45),
             size=Pt(26), bold=True, color=WHITE)
    add_text(sl, title, x + Inches(0.6), y + Inches(0.22), w - Inches(0.65), Inches(0.5),
             size=Pt(15), bold=True, color=WHITE)

    # 본문
    add_rect(sl, x, y + Inches(0.9), w, Inches(4.5),
             fill=LGRAY, line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(1))
    cy2 = y + Inches(1.05)
    for b in bullets:
        prefix = "▸" if not b.startswith("(") else " "
        color  = DGRAY if not b.startswith("(") else GRAY
        add_text(sl, f"{prefix}  {b}", x + Inches(0.15), cy2,
                 w - Inches(0.25), Inches(0.38),
                 size=Pt(11.5), color=color)
        cy2 += Inches(0.38)

    # 화살표
    if i < 3:
        ax = x + w + Inches(0.06)
        add_rect(sl, ax, y + Inches(0.4), Inches(0.1), Inches(0.12), fill=GRAY)

# 하단 보충
add_rect(sl, Inches(0.25), Inches(6.1), W - Inches(0.5), Inches(0.65),
         fill=RGBColor(0xe0, 0xf7, 0xf4), line=TEAL, line_w=Pt(1))
add_text(sl,
         "💡  Semantic Cache: 유사 질문(TF-IDF ≥ 0.75) 감지 시 S1~S4를 건너뛰고 즉시 반환 → LLM 비용 최소화",
         Inches(0.45), Inches(6.15), W - Inches(0.8), Inches(0.55),
         size=Pt(12), color=TEAL)

footer(sl, 4, TOTAL)


# ────────────────────────────────────────────────────────────────
# 5. Web UI 구성
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "Web UI 구성", "브라우저 기반 통합 관리 화면")

# 사이드바 시뮬레이션
add_rect(sl, Inches(0.25), Inches(1.25), Inches(2.4), Inches(5.9), fill=NAVY)
add_text(sl, "pgxllm", Inches(0.38), Inches(1.35), Inches(2.1), Inches(0.5),
         size=Pt(16), bold=True, color=WHITE)
add_text(sl, "Query Test UI", Inches(0.38), Inches(1.78), Inches(2.1), Inches(0.35),
         size=Pt(10), color=TEAL)

nav = [
    ("QUERY",    None,    False),
    ("▶ SQL 실행", True,  False),
    ("CATALOG",  None,    False),
    ("🗂 Schema 탐색", True, False),
    ("GRAPH",    None,    False),
    ("🔗 Graph 관계", True, False),
    ("  📋 Query 분석", True, True),
    ("  🔧 Query 튜닝", True, True),
    ("RULES",    None,    False),
    ("📐 Dialect Rules", True, False),
    ("ADMIN",    None,    False),
    ("🗄 DB 관리",  True,  False),
    ("🤖 LLM 설정", True, False),
]
cy = Inches(2.2)
for label, is_link, is_sub in nav:
    if not is_link:
        add_text(sl, label, Inches(0.42), cy, Inches(2.1), Inches(0.27),
                 size=Pt(8), color=RGBColor(0x9c, 0xa3, 0xaf))
    else:
        indent = Inches(0.55) if is_sub else Inches(0.42)
        add_text(sl, label, indent, cy, Inches(2.0), Inches(0.33),
                 size=Pt(10.5), color=RGBColor(0xd1, 0xd5, 0xdb))
    cy += Inches(0.28) if not is_link else Inches(0.33)

# 페이지별 설명 카드 (우측)
pages_info = [
    ("SQL 실행",      TEAL, "직접 SQL + LLM 자연어 질문\n캐시 히트 시 즉시 반환, 이력 관리"),
    ("Schema 탐색",   NAVY, "테이블·컬럼·인덱스 검색\n↻ 재수집 버튼으로 즉시 갱신"),
    ("Query 분석",    TEAL, "pg_stat 수집 → EXPLAIN 시각화\n독점 수행시간·열색상 표시"),
    ("Query 튜닝",    NAVY, "LLM 최적화 제안\nPlan/ANALYZE 전후 비교"),
    ("Dialect Rules", TEAL, "컬럼별 SQL 작성 규칙\n자동 감지 + 수동 등록"),
    ("LLM 설정",      NAVY, "Provider 실시간 전환\nAPI Key 저장·테스트"),
]
cw = Inches(3.35)
ch = Inches(1.65)
for i, (title, col, body) in enumerate(pages_info):
    r = i // 3
    c = i % 3
    x = Inches(2.85) + c * (cw + Inches(0.12))
    y = Inches(1.25) + r * (ch + Inches(0.12))
    add_rect(sl, x, y, cw, Inches(0.45), fill=col)
    add_text(sl, title, x + Inches(0.12), y + Inches(0.07),
             cw - Inches(0.15), Inches(0.35),
             size=Pt(13), bold=True, color=WHITE)
    add_rect(sl, x, y + Inches(0.45), cw, ch - Inches(0.45),
             fill=LGRAY, line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(0.75))
    add_text(sl, body, x + Inches(0.12), y + Inches(0.52),
             cw - Inches(0.2), ch - Inches(0.6),
             size=Pt(11), color=DGRAY)

# Topbar 배지 설명
add_rect(sl, Inches(2.85), Inches(6.55), Inches(10.2), Inches(0.55),
         fill=RGBColor(0xe6, 0xfb, 0xf7), line=TEAL, line_w=Pt(1))
add_text(sl, "💡  Topbar: 모든 페이지에서 '🤖 Provider · 모델명' 배지로 현재 활성 LLM을 항상 확인",
         Inches(3.0), Inches(6.6), Inches(9.8), Inches(0.45),
         size=Pt(11.5), color=TEAL)

footer(sl, 5, TOTAL)


# ────────────────────────────────────────────────────────────────
# 6. 쿼리 분석 — 실행계획 시각화
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "쿼리 분석 — 실행계획 시각화",
             "EXPLAIN ANALYZE 결과를 방향성 그래프로 표시")

# 좌측: 실행계획 그래프 시뮬레이션
add_text(sl, "실행계획 그래프 (예시)", Inches(0.3), Inches(1.25), Inches(6.0), Inches(0.4),
         size=Pt(13), bold=True, color=NAVY)

plan_nodes = [
    ("Aggregate",              Inches(2.5),  Inches(1.75), GREEN,  "5%"),
    ("Hash Join",              Inches(2.5),  Inches(2.65), ORANGE, "38%"),
    ("Seq Scan\n(orders)",     Inches(0.7),  Inches(3.55), GREEN,  "8%"),
    ("Hash\n(customers)",      Inches(4.3),  Inches(3.55), RED,    "49%"),
]
nw, nh = Inches(1.9), Inches(0.75)
for label, nx, ny, col, pct in plan_nodes:
    add_rect(sl, nx, ny, nw, nh, fill=col)
    add_text(sl, label, nx, ny, nw, nh * 0.65,
             size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, f"exclusive: {pct}", nx, ny + nh * 0.6, nw, nh * 0.4,
             size=Pt(9.5), color=WHITE, align=PP_ALIGN.CENTER)

# 연결선 (단순 직선으로 표현)
for x1, y1, x2, y2 in [
    (Inches(3.45), Inches(2.65), Inches(3.45), Inches(2.5)),
    (Inches(1.65), Inches(3.55), Inches(2.9),  Inches(3.4)),
    (Inches(5.25), Inches(3.55), Inches(4.0),  Inches(3.4)),
]:
    add_rect(sl, x1, y2, Inches(0.06), y1 - y2, fill=GRAY)

# 범례
add_text(sl, "열 색상 범례 (Exclusive 비중)",
         Inches(0.3), Inches(4.75), Inches(5.8), Inches(0.38),
         size=Pt(12), bold=True, color=NAVY)
legend = [
    (GREEN,  "< 10%   정상"),
    (RGBColor(0xfb, 0xbf, 0x24), "10~25%  주의"),
    (ORANGE, "25~50%  병목 가능"),
    (RED,    "≥ 50%   주요 병목"),
]
lx = Inches(0.3)
for col, label in legend:
    add_rect(sl, lx, Inches(5.2), Inches(0.35), Inches(0.3), fill=col)
    add_text(sl, label, lx + Inches(0.42), Inches(5.2), Inches(1.5), Inches(0.3),
             size=Pt(11), color=DGRAY)
    lx += Inches(2.0)

# 우측: $N 파라미터 처리
add_text(sl, "$N 파라미터 처리", Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.4),
         size=Pt(13), bold=True, color=NAVY)

add_rect(sl, Inches(7.0), Inches(1.72), Inches(6.0), Inches(1.6),
         fill=LGRAY, line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(1))
add_text(sl, "pg_stat_statements 추상화:",
         Inches(7.15), Inches(1.82), Inches(5.7), Inches(0.35),
         size=Pt(11), bold=True, color=NAVY)
add_text(sl, 'SELECT * FROM orders\nWHERE status = $1 AND date > $2',
         Inches(7.15), Inches(2.18), Inches(5.7), Inches(0.8),
         size=Pt(11), color=DGRAY)

bullets_param = [
    "$N → 에디터에서 빨간색 강조 표시",
    "파라미터 입력창에 실제 값 입력",
    "Plan/ANALYZE 실행 시 자동 치환",
    "타입 자동 감지: 숫자/bool/null/문자열",
]
cy = Inches(3.5)
for b in bullets_param:
    add_text(sl, f"▸  {b}", Inches(7.0), cy, Inches(6.0), Inches(0.35),
             size=Pt(11.5), color=DGRAY)
    cy += Inches(0.4)

# 노드 조건 아이콘
add_text(sl, "노드 조건 아이콘", Inches(7.0), Inches(5.1), Inches(5.8), Inches(0.38),
         size=Pt(12), bold=True, color=NAVY)
icons = [
    ("🔑 Index Cond",   "⊕ Hash/Merge Cond"),
    ("▽ Filter",        "↕ Sort Key"),
    ("⚡ Join Filter",   "⊞ Group Key"),
]
cy = Inches(5.52)
for left, right in icons:
    add_text(sl, left,  Inches(7.0),  cy, Inches(2.8), Inches(0.32), size=Pt(11), color=DGRAY)
    add_text(sl, right, Inches(10.0), cy, Inches(2.8), Inches(0.32), size=Pt(11), color=DGRAY)
    cy += Inches(0.34)

footer(sl, 6, TOTAL)


# ────────────────────────────────────────────────────────────────
# 7. LLM Provider 지원
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "LLM Provider 지원", "6가지 LLM 서비스 통합 — Web UI에서 실시간 전환")

providers = [
    ("Ollama",           "로컬",   TEAL,
     ["http://localhost:11434", "qwen2.5-coder:7b / 32b", "인증 불필요", "오프라인 운영 가능"]),
    ("vLLM",             "로컬/서버", NAVY,
     ["OpenAI 호환 엔드포인트", "http://localhost:8001/v1", "API Key 선택", "자체 모델 배포"]),
    ("LM Studio",        "로컬",   TEAL,
     ["OpenAI 호환", "http://localhost:1234/v1", "GUI로 모델 관리", "인증 불필요"]),
    ("OpenAI",           "클라우드", NAVY,
     ["api.openai.com", "gpt-4o / gpt-4-turbo", "OPENAI_API_KEY", "또는 UI 직접 입력"]),
    ("Anthropic Claude", "클라우드", TEAL,
     ["api.anthropic.com", "claude-3-5-sonnet", "ANTHROPIC_API_KEY", "또는 UI 직접 입력"]),
    ("IBM watsonx.ai",   "클라우드", NAVY,
     ["us-south / eu-de / jp-tok", "ibm/granite-34b-code", "IBM API Key 필요", "Project ID 필요"]),
]

cw = Inches(2.1)
ch = Inches(4.3)
for i, (name, env, col, items) in enumerate(providers):
    x = Inches(0.25) + i * (cw + Inches(0.09))
    y = Inches(1.25)
    # 헤더
    add_rect(sl, x, y, cw, Inches(0.8), fill=col)
    add_text(sl, name, x + Inches(0.1), y + Inches(0.06), cw - Inches(0.12), Inches(0.45),
             size=Pt(13), bold=True, color=WHITE)
    add_text(sl, env, x + Inches(0.1), y + Inches(0.5), cw - Inches(0.12), Inches(0.28),
             size=Pt(10), color=RGBColor(0xb2, 0xf5, 0xea) if col==TEAL else RGBColor(0x93, 0xc5, 0xfd))
    # 본문
    add_rect(sl, x, y + Inches(0.8), cw, ch - Inches(0.8),
             fill=LGRAY, line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(0.75))
    cy = y + Inches(0.95)
    for item in items:
        add_text(sl, f"▸ {item}", x + Inches(0.1), cy, cw - Inches(0.15), Inches(0.38),
                 size=Pt(10), color=DGRAY)
        cy += Inches(0.38)

# 하단: watsonx 지역
add_rect(sl, Inches(0.25), Inches(5.75), W - Inches(0.5), Inches(0.9),
         fill=RGBColor(0xe0, 0xf7, 0xf4), line=TEAL, line_w=Pt(1))
add_text(sl, "IBM watsonx.ai 지역 Endpoint",
         Inches(0.45), Inches(5.82), Inches(3.5), Inches(0.35),
         size=Pt(12), bold=True, color=NAVY)
regions = [
    "🇺🇸 us-south.ml.cloud.ibm.com",
    "🇩🇪 eu-de.ml.cloud.ibm.com",
    "🇬🇧 eu-gb.ml.cloud.ibm.com",
    "🇯🇵 jp-tok.ml.cloud.ibm.com",
    "🇦🇺 au-syd.ml.cloud.ibm.com",
]
rx = Inches(0.45)
for r in regions:
    add_text(sl, r, rx, Inches(6.18), Inches(2.4), Inches(0.35),
             size=Pt(10.5), color=TEAL)
    rx += Inches(2.42)

footer(sl, 7, TOTAL)


# ────────────────────────────────────────────────────────────────
# 8. LLM 설정 화면
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "LLM 설정 — Web UI", "서버 재시작 없이 실시간 LLM 전환")

# 현재 활성 카드 시뮬레이션
add_rect(sl, Inches(0.3), Inches(1.25), W - Inches(0.6), Inches(1.3),
         fill=WHITE, line=TEAL, line_w=Pt(3))
add_rect(sl, Inches(0.3), Inches(1.25), Inches(0.08), Inches(1.3), fill=TEAL)
add_text(sl, "현재 활성 LLM",
         Inches(0.55), Inches(1.3), Inches(4), Inches(0.3),
         size=Pt(10), bold=True, color=GRAY)
info_items = [
    ("Provider",    "IBM watsonx.ai",                    TEAL),
    ("Model",       "ibm/granite-34b-code-instruct",     NAVY),
    ("Endpoint",    "https://us-south.ml.cloud.ibm.com", DGRAY),
    ("API Key",     "✅ 설정됨",                          GREEN),
    ("timeout",     "600s",                              GRAY),
    ("max_tokens",  "2048",                              GRAY),
    ("temp",        "0.0",                               GRAY),
]
ix = Inches(0.55)
for label, val, col in info_items:
    add_text(sl, label, ix, Inches(1.65), Inches(1.2), Inches(0.28),
             size=Pt(9), color=GRAY)
    add_text(sl, val, ix, Inches(1.93), Inches(1.5), Inches(0.38),
             size=Pt(11), bold=True, color=col)
    ix += Inches(1.82)

# Provider 선택 버튼
add_text(sl, "① Provider 선택",
         Inches(0.3), Inches(2.75), Inches(5), Inches(0.35),
         size=Pt(12), bold=True, color=NAVY)
providers_btns = ["Ollama", "vLLM", "LM Studio", "OpenAI", "Anthropic", "IBM watsonx.ai"]
bx = Inches(0.35)
for i, pv in enumerate(providers_btns):
    is_sel = pv == "IBM watsonx.ai"
    add_rect(sl, bx, Inches(3.15), Inches(1.95), Inches(0.48),
             fill=TEAL if is_sel else WHITE,
             line=TEAL, line_w=Pt(1))
    add_text(sl, pv, bx, Inches(3.15), Inches(1.95), Inches(0.48),
             size=Pt(11), bold=is_sel,
             color=WHITE if is_sel else TEAL,
             align=PP_ALIGN.CENTER)
    bx += Inches(2.0)

# 입력 필드
add_text(sl, "② 접속 정보 입력",
         Inches(0.3), Inches(3.78), Inches(5), Inches(0.35),
         size=Pt(12), bold=True, color=NAVY)

fields = [
    ("Endpoint URL",  "https://us-south.ml.cloud.ibm.com", Inches(7.5)),
    ("API Key",       "•••••••••••••••••••••",              Inches(3.6)),
    ("Project ID",    "xxxxxxxx-xxxx-xxxx-xxxx-xxxxxxxxxxxx", Inches(3.6)),
    ("Model ID",      "ibm/granite-34b-code-instruct",      Inches(7.5)),
]
fy = Inches(4.18)
fx_pairs = [(Inches(0.35), Inches(7.6)), (Inches(8.0), Inches(4.6))]
for i, (label, val, fw) in enumerate(fields):
    fx = Inches(0.35) if i % 2 == 0 else Inches(7.0)
    fy_cur = Inches(4.18) if i < 2 else Inches(4.95)
    add_text(sl, label, fx, fy_cur - Inches(0.25), fw, Inches(0.25),
             size=Pt(9.5), bold=True, color=GRAY)
    add_rect(sl, fx, fy_cur, fw, Inches(0.42),
             fill=WHITE, line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(1))
    add_text(sl, val, fx + Inches(0.1), fy_cur + Inches(0.07), fw - Inches(0.15), Inches(0.32),
             size=Pt(11), color=DGRAY)

# 버튼
add_rect(sl, Inches(0.35), Inches(5.7), Inches(2.5), Inches(0.52),
         fill=WHITE, line=TEAL, line_w=Pt(1.5))
add_text(sl, "🔌  연결 테스트", Inches(0.35), Inches(5.7), Inches(2.5), Inches(0.52),
         size=Pt(13), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(3.1), Inches(5.7), Inches(2.0), Inches(0.52), fill=TEAL)
add_text(sl, "💾  저장", Inches(3.1), Inches(5.7), Inches(2.0), Inches(0.52),
         size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "✅ 연결 성공 — 모델: ibm/granite-34b-code-instruct",
         Inches(5.3), Inches(5.78), Inches(7.5), Inches(0.38),
         size=Pt(12), color=GREEN)

footer(sl, 8, TOTAL)


# ────────────────────────────────────────────────────────────────
# 9. Internal DB 구조
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "Internal DB 구조", "pgxllm 메타데이터 저장소 테이블")

table_rows(
    sl,
    headers=["테이블", "용도", "주요 컬럼"],
    rows=[
        ["db_registry",    "등록된 Target DB 목록",       "alias, host, port, dbname, schema_mode"],
        ["schema_catalog", "수집된 스키마 (테이블·컬럼)", "db_alias, table_name, column_name, data_type"],
        ["graph_edges",    "테이블 간 관계 그래프",        "from_table, to_table, relation_type, approved"],
        ["graph_paths",    "사전 계산된 JOIN 경로",        "from_table, to_table, path_json, depth"],
        ["dialect_rules",  "컬럼별 SQL 작성 규칙",         "rule_id, table_name, column_name, instruction"],
        ["sql_patterns",   "Dynamic SQL 패턴",            "name, detect_keywords, template, enabled"],
        ["query_cache",    "Semantic Cache (TF-IDF)",     "question, sql, db_alias, hit_count"],
        ["query_history",  "쿼리 실행 이력",               "db_alias, mode, question, ok, duration_ms"],
        ["llm_settings",   "LLM 설정 (UI 저장 시)",        "provider, model, api_key, project_id"],
    ],
    x=Inches(0.3),
    y=Inches(1.3),
    w=W - Inches(0.6),
    col_widths=[Inches(2.3), Inches(3.5), Inches(6.7)],
)

add_rect(sl, Inches(0.3), Inches(6.25), W - Inches(0.6), Inches(0.55),
         fill=RGBColor(0xe0, 0xf7, 0xf4), line=TEAL, line_w=Pt(1))
add_text(sl,
         "💡  모든 테이블은 pgxllm 스키마 내에 생성됩니다. CREATE TABLE IF NOT EXISTS 패턴으로 앱 시작 시 자동 생성.",
         Inches(0.5), Inches(6.3), W - Inches(0.9), Inches(0.45),
         size=Pt(11.5), color=TEAL)

footer(sl, 9, TOTAL)


# ────────────────────────────────────────────────────────────────
# 10. 설치 및 CLI
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "설치 및 CLI", "빠른 시작 가이드 및 자동화 명령어")

# 좌측: 설치
add_text(sl, "설치 단계", Inches(0.3), Inches(1.25), Inches(6.0), Inches(0.38),
         size=Pt(13), bold=True, color=NAVY)

install_steps = [
    ("1", "환경 설정",   "bash setup.sh\nsource .venv/bin/activate"),
    ("2", "환경 변수",   "cp .env.example .env\n# DB 접속 정보 입력"),
    ("3", "DB 등록",     "pgxllm db register\n  --alias mydb --host localhost ..."),
    ("4", "스키마 수집", "pgxllm db refresh --alias mydb"),
    ("5", "서버 시작",   "make serve  # 프로덕션 (포트 8000)\nmake dev-backend + dev-frontend"),
]
cy = Inches(1.72)
for num, title, cmd in install_steps:
    add_rect(sl, Inches(0.3), cy, Inches(0.45), Inches(0.45),
             fill=TEAL)
    add_text(sl, num, Inches(0.3), cy, Inches(0.45), Inches(0.45),
             size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, Inches(0.85), cy, Inches(5.2), Inches(0.3),
             size=Pt(11.5), bold=True, color=NAVY)
    add_text(sl, cmd, Inches(0.85), cy + Inches(0.3), Inches(5.2), Inches(0.55),
             size=Pt(10), color=GRAY)
    cy += Inches(1.0)

# 우측: CLI 명령어
add_text(sl, "CLI 명령어 (자동화·설정 전용)", Inches(6.8), Inches(1.25), Inches(6.2), Inches(0.38),
         size=Pt(13), bold=True, color=NAVY)

cli_cmds = [
    ("db register",  "Target DB 등록",
     "--alias mydb --host db.example.com\n--user postgres --dbname mydb"),
    ("db refresh",   "스키마 재수집 (cron 자동화)",
     "--alias mydb  또는  --all\n--skip-samples --skip-rules"),
    ("web",          "Web UI 서버 시작",
     "--host 0.0.0.0 --port 8000\n--reload (개발용)"),
    ("eval",         "BIRD Benchmark 평가",
     "--file bird_dev.json --alias mydb\n--output results/eval.json"),
]
cy = Inches(1.72)
for cmd, desc, args in cli_cmds:
    add_rect(sl, Inches(6.8), cy, Inches(6.2), Inches(1.15),
             fill=LGRAY, line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(1))
    add_text(sl, f"pgxllm {cmd}", Inches(6.95), cy + Inches(0.1),
             Inches(5.8), Inches(0.38), size=Pt(12), bold=True, color=NAVY)
    add_text(sl, desc, Inches(6.95), cy + Inches(0.45),
             Inches(5.8), Inches(0.3), size=Pt(10.5), color=TEAL)
    add_text(sl, args, Inches(6.95), cy + Inches(0.75),
             Inches(5.8), Inches(0.38), size=Pt(10), color=GRAY)
    cy += Inches(1.25)

footer(sl, 10, TOTAL)


# ────────────────────────────────────────────────────────────────
# 11. REST API 레퍼런스
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
slide_header(sl, "REST API 레퍼런스", "Base URL: http://서버:8000/api")

api_groups = [
    ("DB 관리",     TEAL, [
        ("GET",    "/db/list",              "등록된 DB 목록"),
        ("POST",   "/db/register",          "DB 등록"),
        ("POST",   "/db/refresh/{alias}",   "스키마 재수집"),
        ("DELETE", "/db/{alias}",           "DB 삭제"),
    ]),
    ("쿼리",        NAVY, [
        ("POST",   "/query/run",            "SQL 직접 실행 / LLM 파이프라인"),
        ("GET",    "/query/history",        "실행 이력"),
        ("DELETE", "/query/cache/all",      "전체 캐시 삭제"),
    ]),
    ("pg_stat 분석", TEAL, [
        ("GET",    "/pgstat/{alias}/queries",       "수집된 쿼리 목록"),
        ("POST",   "/pgstat/{alias}/query/tune",    "LLM 튜닝 제안"),
        ("POST",   "/pgstat/{alias}/query/plan",    "EXPLAIN 실행계획"),
        ("POST",   "/pgstat/{alias}/query/describe","LLM 쿼리 설명"),
    ]),
    ("Graph",       NAVY, [
        ("GET",    "/graph/{alias}",                "관계 엣지 목록"),
        ("POST",   "/graph/{alias}/approve-all",    "전체 승인"),
        ("POST",   "/graph/{alias}/refresh-paths",  "경로 재계산"),
    ]),
    ("LLM 설정",    TEAL, [
        ("GET",    "/llm/providers",   "지원 Provider 목록"),
        ("GET",    "/llm/config",      "현재 설정 (API Key 마스킹)"),
        ("POST",   "/llm/config",      "설정 저장"),
        ("POST",   "/llm/test",        "연결 테스트"),
    ]),
    ("스키마 / Rules", NAVY, [
        ("GET",    "/schema/{alias}",         "테이블·컬럼 목록 (?search=)"),
        ("GET",    "/schema/{alias}/indexes", "인덱스 목록"),
        ("GET",    "/rules/{alias}",          "Dialect Rules 목록"),
        ("POST",   "/rules/{alias}",          "규칙 추가"),
    ]),
]

cw = Inches(2.1)
ch = Inches(2.85)
for i, (group, col, endpoints) in enumerate(api_groups):
    r = i // 3
    c = i % 3
    x = Inches(0.25) + c * (cw + Inches(0.13))
    y = Inches(1.25) + r * (ch + Inches(0.12))
    add_rect(sl, x, y, cw, Inches(0.4), fill=col)
    add_text(sl, group, x + Inches(0.1), y + Inches(0.06),
             cw - Inches(0.15), Inches(0.32), size=Pt(12), bold=True, color=WHITE)
    add_rect(sl, x, y + Inches(0.4), cw, ch - Inches(0.4),
             fill=LGRAY, line=RGBColor(0xd1, 0xd5, 0xdb), line_w=Pt(0.75))
    cy = y + Inches(0.52)
    for method, path, desc in endpoints:
        mcol = {"GET": GREEN, "POST": TEAL, "DELETE": RED, "PATCH": ORANGE}[method]
        add_rect(sl, x + Inches(0.1), cy, Inches(0.55), Inches(0.23), fill=mcol)
        add_text(sl, method, x + Inches(0.1), cy, Inches(0.55), Inches(0.23),
                 size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, path, x + Inches(0.7), cy, cw - Inches(0.75), Inches(0.23),
                 size=Pt(8.5), color=NAVY)
        add_text(sl, desc, x + Inches(0.12), cy + Inches(0.24), cw - Inches(0.18), Inches(0.25),
                 size=Pt(8.5), color=GRAY)
        cy += Inches(0.55)

footer(sl, 11, TOTAL)


# ────────────────────────────────────────────────────────────────
# 12. 마무리 / 로드맵
# ────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, H - Inches(1.5), W, Inches(1.5), fill=TEAL)

add_text(sl, "pgxllm",
         Inches(1.0), Inches(1.5), Inches(11), Inches(1.2),
         size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "PostgreSQL Text-to-SQL System",
         Inches(1.0), Inches(2.7), Inches(11), Inches(0.6),
         size=Pt(20), color=TEAL, align=PP_ALIGN.CENTER)

# 요약 포인트
summary = [
    "4단계 파이프라인으로 자연어 → 정확한 SQL 변환",
    "6가지 LLM Provider 지원 — 로컬부터 IBM watsonx.ai까지",
    "Web UI에서 스키마·쿼리·관계·규칙 통합 관리",
    "EXPLAIN ANALYZE 실행계획 시각화 + 독점 수행시간 분석",
    "TF-IDF Semantic Cache로 반복 질문 즉시 응답",
]
cy = Inches(3.6)
for s in summary:
    add_text(sl, f"✓   {s}",
             Inches(2.5), cy, Inches(8.5), Inches(0.42),
             size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
    cy += Inches(0.46)

add_text(sl, "자연어로 묻고, SQL로 답하다",
         Inches(1.0), H - Inches(1.2), Inches(11), Inches(0.6),
         size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(sl, 12, TOTAL)


# ── 저장 ──────────────────────────────────────────────────────
import os
os.makedirs("docs", exist_ok=True)
out = "docs/pgxllm.pptx"
prs.save(out)
print(f"✅ 저장 완료: {out}")
print(f"   슬라이드 수: {TOTAL}")
